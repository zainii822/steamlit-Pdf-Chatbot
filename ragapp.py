import os
import io
import hashlib
from typing import List, Dict

import streamlit as st
from openai import OpenAI

from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from PIL import Image
import pytesseract


# ============================================================
# PAGE CONFIG
# ============================================================

st.set_page_config(
    page_title="Document Chatbot",
    page_icon="📚",
    layout="wide"
)


# ============================================================
# CONSTANTS
# ============================================================

ALLOWED_EXTENSIONS = {
    "pdf",
    "docx",
    "xlsx",
    "pptx"
}

DEFAULT_MODEL = "gpt-4o-mini"
DEFAULT_EMBEDDING_MODEL = "text-embedding-3-small"


# ============================================================
# SESSION STATE
# ============================================================

if "documents" not in st.session_state:
    st.session_state.documents = []

if "chunks" not in st.session_state:
    st.session_state.chunks = []

if "messages" not in st.session_state:
    st.session_state.messages = []

if "file_hashes" not in st.session_state:
    st.session_state.file_hashes = set()


# ============================================================
# SIDEBAR
# ============================================================

with st.sidebar:
    st.title("⚙️ Settings")

    api_key = st.text_input(
        "OpenAI API Key",
        type="password",
        placeholder="sk-..."
    )

    st.divider()

    st.subheader("Chunk Settings")

    chunk_size = st.number_input(
        "Chunk Size",
        min_value=200,
        max_value=5000,
        value=1000,
        step=100
    )

    chunk_overlap = st.number_input(
        "Chunk Overlap",
        min_value=0,
        max_value=1000,
        value=150,
        step=50
    )

    if chunk_overlap >= chunk_size:
        st.error("Chunk overlap must be smaller than chunk size.")

    st.divider()

    model = st.selectbox(
        "Chat Model",
        [
            "gpt-4o-mini",
            "gpt-4o"
        ],
        index=0
    )

    st.divider()

    if st.button("🗑️ Clear Everything", use_container_width=True):
        st.session_state.documents = []
        st.session_state.chunks = []
        st.session_state.messages = []
        st.session_state.file_hashes = set()

        st.rerun()


# ============================================================
# TITLE
# ============================================================

st.title("📚 Document Chatbot")

st.markdown(
    """
    Upload **PDF, DOCX, XLSX or PPTX** files and ask questions
    only about their content.

    **Unsupported files → Invalid file**  
    **Information not present in documents → Not found**
    """
)


# ============================================================
# VALIDATE FILE
# ============================================================

def get_extension(filename: str) -> str:
    if "." not in filename:
        return ""

    return filename.rsplit(".", 1)[1].lower()


def validate_file(uploaded_file):
    extension = get_extension(uploaded_file.name)

    if extension not in ALLOWED_EXTENSIONS:
        return False

    return True


# ============================================================
# IMAGE OCR
# ============================================================

def extract_image_text(image_bytes: bytes) -> str:
    """
    Extract text from image using OCR.
    """

    try:
        image = Image.open(io.BytesIO(image_bytes))

        text = pytesseract.image_to_string(
            image,
            lang="eng"
        )

        return text.strip()

    except Exception:
        return ""


# ============================================================
# PDF
# ============================================================

def extract_pdf(file_bytes: bytes) -> str:

    text_parts = []

    reader = PdfReader(
        io.BytesIO(file_bytes)
    )

    for page_number, page in enumerate(reader.pages, start=1):

        page_text = page.extract_text() or ""

        if page_text.strip():
            text_parts.append(
                f"\n[PDF Page {page_number}]\n"
                f"{page_text}"
            )

        # ----------------------------------------------------
        # Images
        # ----------------------------------------------------

        try:
            for image_index, image in enumerate(
                page.images,
                start=1
            ):

                image_bytes = image.data

                image_text = extract_image_text(
                    image_bytes
                )

                if image_text:
                    text_parts.append(
                        f"\n[PDF Page {page_number} "
                        f"Image {image_index}]\n"
                        f"{image_text}"
                    )

        except Exception:
            pass

    return "\n".join(text_parts)


# ============================================================
# DOCX
# ============================================================

def extract_docx(file_bytes: bytes) -> str:

    text_parts = []

    doc = Document(
        io.BytesIO(file_bytes)
    )

    # Paragraphs
    for paragraph in doc.paragraphs:

        text = paragraph.text.strip()

        if text:
            text_parts.append(text)

    # Tables
    for table in doc.tables:

        for row in table.rows:

            row_text = []

            for cell in row.cells:
                row_text.append(
                    cell.text.strip()
                )

            text_parts.append(
                " | ".join(row_text)
            )

    # Images
    for relation in doc.part.rels.values():

        if "image" in relation.reltype:

            try:

                image_bytes = relation.target_part.blob

                image_text = extract_image_text(
                    image_bytes
                )

                if image_text:
                    text_parts.append(
                        "\n[DOCX Image]\n"
                        + image_text
                    )

            except Exception:
                pass

    return "\n".join(text_parts)


# ============================================================
# XLSX
# ============================================================

def extract_xlsx(file_bytes: bytes) -> str:

    text_parts = []

    workbook = load_workbook(
        filename=io.BytesIO(file_bytes),
        data_only=True
    )

    for sheet in workbook.worksheets:

        text_parts.append(
            f"\n[Excel Sheet: {sheet.title}]"
        )

        for row in sheet.iter_rows(
            values_only=True
        ):

            values = []

            for value in row:

                if value is not None:
                    values.append(str(value))

            if values:

                text_parts.append(
                    " | ".join(values)
                )

    return "\n".join(text_parts)


# ============================================================
# PPTX
# ============================================================

def extract_pptx(file_bytes: bytes) -> str:

    text_parts = []

    presentation = Presentation(
        io.BytesIO(file_bytes)
    )

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        text_parts.append(
            f"\n[PPTX Slide {slide_number}]"
        )

        for shape in slide.shapes:

            # Text
            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text:
                    text_parts.append(text)

            # Images
            if shape.shape_type == 13:

                try:

                    image_bytes = (
                        shape.image.blob
                    )

                    image_text = extract_image_text(
                        image_bytes
                    )

                    if image_text:
                        text_parts.append(
                            f"\n[Slide {slide_number} "
                            f"Image]\n"
                            f"{image_text}"
                        )

                except Exception:
                    pass

    return "\n".join(text_parts)


# ============================================================
# GENERIC EXTRACTION
# ============================================================

def extract_document(file_name, file_bytes):

    extension = get_extension(file_name)

    if extension == "pdf":
        return extract_pdf(file_bytes)

    if extension == "docx":
        return extract_docx(file_bytes)

    if extension == "xlsx":
        return extract_xlsx(file_bytes)

    if extension == "pptx":
        return extract_pptx(file_bytes)

    return None


# ============================================================
# CHUNKING
# ============================================================

def create_chunks(
    text: str,
    chunk_size: int,
    overlap: int
) -> List[str]:

    if not text:
        return []

    text = text.strip()

    chunks = []

    start = 0
    text_length = len(text)

    while start < text_length:

        end = min(
            start + chunk_size,
            text_length
        )

        chunk = text[start:end].strip()

        if chunk:
            chunks.append(chunk)

        if end >= text_length:
            break

        start = end - overlap

    return chunks


# ============================================================
# EMBEDDINGS
# ============================================================

def get_embeddings(
    client: OpenAI,
    texts: List[str]
):

    response = client.embeddings.create(
        model=DEFAULT_EMBEDDING_MODEL,
        input=texts
    )

    return [
        item.embedding
        for item in response.data
    ]


# ============================================================
# COSINE SIMILARITY
# ============================================================

def cosine_similarity(a, b):

    dot_product = sum(
        x * y
        for x, y in zip(a, b)
    )

    norm_a = sum(
        x * x
        for x in a
    ) ** 0.5

    norm_b = sum(
        y * y
        for y in b
    ) ** 0.5

    if norm_a == 0 or norm_b == 0:
        return 0

    return dot_product / (
        norm_a * norm_b
    )


# ============================================================
# BUILD INDEX
# ============================================================

def build_index(
    client: OpenAI,
    documents,
    chunk_size,
    chunk_overlap
):

    all_chunks = []

    for document in documents:

        chunks = create_chunks(
            document["text"],
            chunk_size,
            chunk_overlap
        )

        for index, chunk in enumerate(chunks):

            all_chunks.append({
                "text": chunk,
                "file": document["name"],
                "chunk": index
            })

    if not all_chunks:
        return []

    texts = [
        item["text"]
        for item in all_chunks
    ]

    embeddings = get_embeddings(
        client,
        texts
    )

    for item, embedding in zip(
        all_chunks,
        embeddings
    ):
        item["embedding"] = embedding

    return all_chunks


# ============================================================
# SEARCH
# ============================================================

def search_documents(
    client,
    chunks,
    question,
    top_k=5
):

    if not chunks:
        return []

    question_embedding = get_embeddings(
        client,
        [question]
    )[0]

    scored = []

    for item in chunks:

        score = cosine_similarity(
            question_embedding,
            item["embedding"]
        )

        scored.append(
            (score, item)
        )

    scored.sort(
        key=lambda x: x[0],
        reverse=True
    )

    return scored[:top_k]


# ============================================================
# ASK LLM
# ============================================================

def answer_question(
    client,
    model,
    question,
    search_results
):

    if not search_results:
        return "Not found"

    # --------------------------------------------------------
    # Minimum similarity threshold
    # --------------------------------------------------------

    best_score = search_results[0][0]

    # This threshold prevents unrelated questions
    # from receiving a general AI answer.
    if best_score < 0.25:
        return "Not found"

    context_parts = []

    for score, item in search_results:

        context_parts.append(
            f"""
FILE: {item['file']}
RELEVANCE: {score:.3f}

CONTENT:
{item['text']}
"""
        )

    context = "\n".join(
        context_parts
    )

    system_prompt = """
You are a strict document question-answering chatbot.

Your ONLY source of information is the provided DOCUMENT CONTEXT.

Rules:

1. Answer ONLY using the document context.
2. Do not use your own general knowledge.
3. Do not guess.
4. Do not invent information.
5. If the answer is not clearly present in the context, reply exactly:
   Not found
6. If the user asks something unrelated to the uploaded documents,
   reply exactly:
   Not found
7. If the context contains conflicting information, mention the
   conflict and identify the relevant file/page/slide when possible.
8. Keep answers concise but useful.
"""

    user_prompt = f"""
DOCUMENT CONTEXT:

{context}

USER QUESTION:

{question}

Answer strictly from the document context.
"""

    response = client.chat.completions.create(
        model=model,
        temperature=0,
        messages=[
            {
                "role": "system",
                "content": system_prompt
            },
            {
                "role": "user",
                "content": user_prompt
            }
        ]
    )

    answer = response.choices[0].message.content.strip()

    return answer


# ============================================================
# FILE UPLOADER
# ============================================================

uploaded_files = st.file_uploader(
    "📁 Upload Documents",
    type=[
        "pdf",
        "docx",
        "xlsx",
        "pptx"
    ],
    accept_multiple_files=True
)


# ============================================================
# PROCESS FILES
# ============================================================

if uploaded_files:

    if not api_key:

        st.warning(
            "Please enter your OpenAI API key "
            "from the sidebar."
        )

    elif chunk_overlap >= chunk_size:

        st.error(
            "Chunk overlap must be smaller "
            "than chunk size."
        )

    else:

        client = OpenAI(
            api_key=api_key
        )

        new_documents = []

        for uploaded_file in uploaded_files:

            file_bytes = uploaded_file.getvalue()

            file_hash = hashlib.md5(
                file_bytes
            ).hexdigest()

            # ------------------------------------------------
            # Invalid file check
            # ------------------------------------------------

            if not validate_file(uploaded_file):

                st.error(
                    f"{uploaded_file.name}: Invalid file"
                )

                continue

            # ------------------------------------------------
            # Avoid duplicate files
            # ------------------------------------------------

            if file_hash in st.session_state.file_hashes:
                continue

            with st.spinner(
                f"Reading {uploaded_file.name}..."
            ):

                try:

                    extracted_text = extract_document(
                        uploaded_file.name,
                        file_bytes
                    )

                    if not extracted_text.strip():

                        st.warning(
                            f"{uploaded_file.name}: "
                            "No readable text found."
                        )

                        continue

                    document = {
                        "name": uploaded_file.name,
                        "text": extracted_text,
                        "hash": file_hash
                    }

                    new_documents.append(
                        document
                    )

                    st.session_state.file_hashes.add(
                        file_hash
                    )

                except Exception as e:

                    st.error(
                        f"Error reading "
                        f"{uploaded_file.name}: {e}"
                    )

        if new_documents:

            st.session_state.documents.extend(
                new_documents
            )

            with st.spinner(
                "Creating document index..."
            ):

                st.session_state.chunks = build_index(
                    client,
                    st.session_state.documents,
                    chunk_size,
                    chunk_overlap
                )

            st.success(
                f"{len(new_documents)} document(s) "
                "processed successfully."
            )


# ============================================================
# SHOW DOCUMENTS
# ============================================================

if st.session_state.documents:

    st.subheader("📄 Uploaded Documents")

    for document in st.session_state.documents:

        st.write(
            f"✅ {document['name']}"
        )

    st.caption(
        f"Total chunks: "
        f"{len(st.session_state.chunks)}"
    )


# ============================================================
# CHAT HISTORY
# ============================================================

for message in st.session_state.messages:

    with st.chat_message(
        message["role"]
    ):
        st.markdown(
            message["content"]
        )


# ============================================================
# CHAT INPUT
# ============================================================

question = st.chat_input(
    "Ask something about your documents..."
)


if question:

    # --------------------------------------------------------
    # User message
    # --------------------------------------------------------

    st.session_state.messages.append({
        "role": "user",
        "content": question
    })

    with st.chat_message("user"):
        st.markdown(question)

    # --------------------------------------------------------
    # Validation
    # --------------------------------------------------------

    if not api_key:

        answer = "Please enter your API key."

    elif not st.session_state.documents:

        answer = "Not found"

    else:

        client = OpenAI(
            api_key=api_key
        )

        with st.chat_message("assistant"):

            with st.spinner("Searching documents..."):

                results = search_documents(
                    client,
                    st.session_state.chunks,
                    question,
                    top_k=5
                )

                answer = answer_question(
                    client,
                    model,
                    question,
                    results
                )

            st.markdown(answer)

    # --------------------------------------------------------
    # Save assistant response
    # --------------------------------------------------------

    st.session_state.messages.append({
        "role": "assistant",
        "content": answer
    })
